"""品牌合规审核平台 - 文档解析服务"""

import io
import json
import logging
import re
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation

from src.models.schemas import (
    BrandRules,
    ColorRule,
    ColorRules,
    CopywritingRules,
    ForbiddenWord,
    FontRules,
    LayoutRules,
    LogoRules,
    SecondaryRule,
)
from src.utils.config import settings
from src.utils.json_parser import parse_json_response

logger = logging.getLogger(__name__)

# 品牌规范解析 Prompt（DeepSeek 文本模型）
PARSE_SYSTEM_PROMPT = """你是品牌规范文档解析专家。请从品牌规范文档中提取结构化规则信息。

【规则分类】共两类，严格区分：

一、【前置条件 preconditions】- 审核前需用户填写的表单字段（如"品牌标识情况"、"传播类型"、"物料类型"等）。
这类内容绝对不能放入 rules，必须单独放入 preconditions 数组。

二、【审核规则 rules】- 所有审核规则，无论是 Logo、色彩、字体、品牌调性、排版布局还是场景适配，全部按原文编号逐条放入 rules 数组，不得遗漏。
用 category 字段标记分类，按以下顺序输出：Logo → 色彩 → 字体 → 品牌调性 → 排版布局 → 其他

【核心要求：三段式判定条件】
每条规则必须包含三个判定条件字段：

- fail_condition：判定为"违规(FAIL)"的具体情况
- review_condition：需要"人工复核(REVIEW)"的情况
- pass_condition：判定为"合规(PASS)"的情况

判断原则：
- 客观可量化的违规 → fail_condition 写具体阈值
- 主观审美类（品牌调性等）→ review_condition 要详细，fail_condition 只写极端情况
- 图中看不清楚、无法确认 → 写入 review_condition

【输出JSON格式】：
```json
{
  "brand_name": "品牌名称",
  "preconditions": [
    {
      "field_name": "品牌标识情况",
      "required": true,
      "type": "单选",
      "options": ["包含（作为常规品牌标识）", "包含（Logo即为画面核心主体）", "无（特殊说明）"],
      "logic": "若选Logo即为画面核心主体，豁免Logo位置和尺寸判定，仅保留颜色与形变校验；若选无，跳过所有Logo规则"
    }
  ],
  "rules": [
    {
      "category": "Logo",
      "name": "品牌Logo是否缺失",
      "content": "正式成稿或正式传播画面中应出现品牌Logo或联合标识",
      "priority": 1,
      "rule_source_id": "H-LOGO-01",
      "fail_condition": "正式传播物料中完全未检测到任何品牌Logo或联合标识",
      "review_condition": "画面用途不明确，无法确认是否属于正式传播物料；Logo区域被遮挡或模糊无法确认",
      "pass_condition": "画面中可清晰识别品牌Logo或联合标识"
    },
    {
      "category": "色彩",
      "name": "主色合规",
      "content": "主色应为#113655（标准蓝版），整体色彩风格健康、明亮、清爽",
      "priority": 1,
      "rule_source_id": "H-COLOR-00",
      "fail_condition": "画面主色明显偏离#113655，且无合理的反白/黑白版本使用场景",
      "review_condition": "图片偏色或压缩导致颜色判断不稳定",
      "pass_condition": "主色符合#113655标准蓝版或规范的反白/黑白版本"
    }
  ]
}
```

严格执行：
1. 所有审核规则统一放入 rules 数组，不得在 rules 数组外单独输出 color/logo/font 结构
2. 前置条件表单字段必须放入 preconditions，绝对不能放入 rules
3. rules 必须按原文编号逐条提取，不得合并，不得遗漏
4. 每条 rules 必须填写 fail_condition、review_condition、pass_condition 三个字段
5. priority: 1=重要(客观可判断), 2=一般(主观判断), 3=参考
6. 只输出JSON，不要其他文字"""

# 支持的文档格式
SUPPORTED_FORMATS = {
    # 格式扩展名 -> (解析方法名, 描述)
    ".pdf": ("_parse_pdf", "PDF文档"),
    ".ppt": ("_parse_ppt", "PowerPoint演示文稿"),
    ".pptx": ("_parse_ppt", "PowerPoint演示文稿"),
    ".doc": ("_parse_word", "Word文档"),
    ".docx": ("_parse_word", "Word文档"),
    ".xls": ("_parse_excel", "Excel表格"),
    ".xlsx": ("_parse_excel", "Excel表格"),
    ".md": ("_parse_text", "Markdown文档"),
    ".txt": ("_parse_text", "文本文件"),
}


class DocumentParser:
    """文档解析服务"""

    def parse(self, file_data: bytes, filename: str) -> BrandRules:
        """解析文档"""
        ext = Path(filename).suffix.lower()

        if ext not in SUPPORTED_FORMATS:
            raise ValueError(f"不支持的文档格式: {ext}，支持的格式: {', '.join(SUPPORTED_FORMATS.keys())}")

        method_name, _ = SUPPORTED_FORMATS[ext]
        parse_method = getattr(self, method_name)
        return parse_method(file_data, filename)

    def parse_file(self, file_path: str, brand_name: str = None, progress_callback=None) -> BrandRules:
        """解析本地文件"""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        with open(file_path, "rb") as f:
            file_data = f.read()

        rules = self.parse(file_data, file_path.name)
        if brand_name:
            rules.brand_name = brand_name

        return rules

    def _parse_pdf(self, file_data: bytes, filename: str) -> BrandRules:
        """解析PDF文档"""
        logger.info(f"解析PDF文档: {filename}")

        doc = fitz.open(stream=file_data, filetype="pdf")
        text_content = []

        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                text_content.append(f"=== 第{page_num + 1}页 ===\n{text}")

        doc.close()

        full_text = "\n\n".join(text_content)
        logger.info(f"PDF解析完成，共{len(text_content)}页，{len(full_text)}字符")

        rules = self._extract_rules_with_llm(full_text, filename)
        rules.raw_text = full_text[:50000]

        return rules

    def _parse_ppt(self, file_data: bytes, filename: str) -> BrandRules:
        """解析PPT/PPTX文档"""
        logger.info(f"解析PPT文档: {filename}")

        prs = Presentation(io.BytesIO(file_data))
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                text_content.append(f"=== 第{slide_num}页 ===\n" + "\n".join(slide_texts))

        full_text = "\n\n".join(text_content)
        logger.info(f"PPT解析完成，共{len(prs.slides)}页，{len(full_text)}字符")

        rules = self._extract_rules_with_llm(full_text, filename)
        rules.raw_text = full_text[:50000]

        return rules

    def _parse_word(self, file_data: bytes, filename: str) -> BrandRules:
        """解析Word文档 (.doc, .docx)"""
        logger.info(f"解析Word文档: {filename}")

        ext = Path(filename).suffix.lower()
        text_content = []

        try:
            if ext == ".docx":
                # 使用python-docx解析.docx
                from docx import Document
                doc = Document(io.BytesIO(file_data))

                # 提取段落文本
                for i, para in enumerate(doc.paragraphs):
                    if para.text.strip():
                        text_content.append(para.text.strip())

                # 提取表格文本
                for table_idx, table in enumerate(doc.tables):
                    table_texts = []
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            table_texts.append(row_text)
                    if table_texts:
                        text_content.append(f"\n=== 表格{table_idx + 1} ===\n" + "\n".join(table_texts))

            else:
                # .doc格式，尝试使用antiword或直接用python-docx（可能失败）
                try:
                    from docx import Document
                    doc = Document(io.BytesIO(file_data))
                    for para in doc.paragraphs:
                        if para.text.strip():
                            text_content.append(para.text.strip())
                except Exception as e:
                    logger.warning(f"解析.doc格式失败，尝试其他方式: {e}")
                    # 尝试作为纯文本解析
                    try:
                        text_content.append(file_data.decode('utf-8', errors='ignore'))
                    except:
                        text_content.append(file_data.decode('gbk', errors='ignore'))

        except ImportError:
            logger.error("未安装python-docx库，无法解析Word文档")
            raise ImportError("请安装python-docx库: pip install python-docx")
        except Exception as e:
            logger.error(f"Word文档解析失败: {e}")
            raise

        full_text = "\n\n".join(text_content)
        logger.info(f"Word解析完成，{len(full_text)}字符")

        rules = self._extract_rules_with_llm(full_text, filename)
        rules.raw_text = full_text[:50000]

        return rules

    def _parse_excel(self, file_data: bytes, filename: str) -> BrandRules:
        """解析Excel表格 (.xls, .xlsx)"""
        logger.info(f"解析Excel文档: {filename}")

        ext = Path(filename).suffix.lower()
        text_content = []

        try:
            if ext == ".xlsx":
                # 使用openpyxl解析.xlsx
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(file_data), data_only=True)

                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_texts = [f"=== 工作表: {sheet_name} ==="]

                    for row in sheet.iter_rows(values_only=True):
                        # 过滤空值，转字符串
                        row_values = [str(cell) if cell is not None else "" for cell in row]
                        row_text = " | ".join(row_values)
                        if row_text.strip(" |"):
                            sheet_texts.append(row_text)

                    if len(sheet_texts) > 1:  # 有内容才添加
                        text_content.append("\n".join(sheet_texts))

            else:
                # .xls格式，使用xlrd
                try:
                    import xlrd
                    workbook = xlrd.open_workbook(file_contents=file_data)

                    for sheet in workbook.sheets():
                        sheet_texts = [f"=== 工作表: {sheet.name} ==="]

                        for row_idx in range(sheet.nrows):
                            row_values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                            row_text = " | ".join(row_values)
                            if row_text.strip(" |"):
                                sheet_texts.append(row_text)

                        if len(sheet_texts) > 1:
                            text_content.append("\n".join(sheet_texts))

                except ImportError:
                    logger.warning("未安装xlrd库，尝试用openpyxl解析.xls")
                    from openpyxl import load_workbook
                    wb = load_workbook(io.BytesIO(file_data), data_only=True)

                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        sheet_texts = [f"=== 工作表: {sheet_name} ==="]

                        for row in sheet.iter_rows(values_only=True):
                            row_values = [str(cell) if cell is not None else "" for cell in row]
                            row_text = " | ".join(row_values)
                            if row_text.strip(" |"):
                                sheet_texts.append(row_text)

                        if len(sheet_texts) > 1:
                            text_content.append("\n".join(sheet_texts))

        except ImportError as e:
            logger.error(f"缺少必要的库: {e}")
            raise ImportError("请安装openpyxl库: pip install openpyxl (或 xlrd for .xls)")
        except Exception as e:
            logger.error(f"Excel解析失败: {e}")
            raise

        full_text = "\n\n".join(text_content)
        logger.info(f"Excel解析完成，{len(full_text)}字符")

        rules = self._extract_rules_with_llm(full_text, filename)
        rules.raw_text = full_text[:50000]

        return rules

    def _parse_text(self, file_data: bytes, filename: str) -> BrandRules:
        """解析文本文件 (.txt, .md)"""
        ext = Path(filename).suffix.lower()
        logger.info(f"解析文本文档: {filename}")

        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
        text_content = None

        for encoding in encodings:
            try:
                text_content = file_data.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text_content is None:
            # 最后尝试忽略错误
            text_content = file_data.decode('utf-8', errors='ignore')

        full_text = text_content.strip()
        logger.info(f"文本解析完成，{len(full_text)}字符")

        # 对于Markdown文件，可以做一些简单的格式处理说明
        if ext == ".md":
            # 统计标题数量
            heading_count = len(re.findall(r'^#{1,6}\s', full_text, re.MULTILINE))
            if heading_count > 0:
                logger.info(f"检测到{heading_count}个Markdown标题")

        rules = self._extract_rules_with_llm(full_text, filename)
        rules.raw_text = full_text[:50000]

        return rules

    def _extract_rules_with_llm(self, text: str, filename: str) -> BrandRules:
        """使用LLM提取品牌规范"""
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI

        logger.info(f"使用LLM提取品牌规范: {filename}")

        # 检查API配置
        if not settings.llm_api_key:
            logger.warning("未配置文本模型 API Key，无法进行LLM提取")
            return BrandRules(
                brand_id="",
                brand_name="",
                version="1.0",
                source=filename,
            )

        system_prompt = PARSE_SYSTEM_PROMPT

        # 截取文本，保留关键内容（提高上限至60K，适应当前大模型64K+上下文窗口）
        max_chars = 60000
        if len(text) > max_chars:
            # 尝试保留前面部分和包含关键词的部分
            text_for_llm = text[:max_chars]
            logger.info(f"文本过长({len(text)}字符)，截取前{max_chars}字符")
        else:
            text_for_llm = text

        user_prompt = f"""请从以下品牌规范文档中提取结构化规则。

文档内容：
{text_for_llm}

请输出JSON格式的规则："""

        try:
            # 使用 DeepSeek 纯文本模型进行规则解析
            logger.info(f"文本模型配置: base={settings.llm_api_base}, model={settings.llm_model}")

            llm = ChatOpenAI(
                model=settings.llm_model,
                base_url=settings.llm_api_base,
                api_key=settings.llm_api_key,
                temperature=0,
                timeout=300,  # 规则文档输出量大，需要更长超时
                max_tokens=16384,  # 规则文档解析输出量大，必须显式设置
            )

            messages = [
                SystemMessage(content=system_prompt),
                HumanMessage(content=user_prompt),
            ]

            logger.info("正在调用LLM（非流式，超时300s）...")
            try:
                response = llm.invoke(messages)
                content = response.content
            except Exception as invoke_err:
                # 非流式超时时回退到流式，逐chunk接收避免超时
                logger.warning(f"非流式调用失败({invoke_err})，回退到流式接收...")
                return self._extract_rules_with_llm_stream(text, filename)

            logger.info(f"LLM响应长度: {len(content)} 字符")
            logger.debug(f"LLM响应内容: {content[:1000]}...")

            # 提取JSON - 多种方式尝试
            data = self._parse_json_response(content)

            if data is None:
                logger.error("无法从LLM响应中解析JSON")
                # 写入完整响应到文件，方便排查（日志终端会截断）
                try:
                    import time
                    dump_path = Path(__file__).parent.parent.parent / "docs" / f"llm_response_dump_{int(time.time())}.txt"
                    dump_path.parent.mkdir(parents=True, exist_ok=True)
                    dump_path.write_text(content, encoding="utf-8")
                    logger.error(f"完整LLM响应已写入: {dump_path}  (共{len(content)}字符)")
                except Exception as dump_err:
                    logger.error(f"写入dump失败: {dump_err}")
                    logger.error(f"原始响应(前500字符): {content[:500]}")
                return BrandRules(
                    brand_id="",
                    brand_name="",
                    version="1.0",
                    source=filename,
                )

            # 构建BrandRules
            rules = BrandRules(
                brand_id="",
                brand_name=data.get("brand_name") or "",
                version="1.0",
                source=filename,
            )

            # 解析色彩（兼容旧格式）
            if data.get("color"):
                self._parse_color_rules(rules, data["color"])

            # 解析Logo（兼容旧格式）
            if data.get("logo"):
                self._parse_logo_rules(rules, data["logo"])

            # 解析字体（兼容旧格式）
            if data.get("font"):
                self._parse_font_rules(rules, data["font"])

            # 解析文案（兼容旧格式）
            if data.get("copywriting"):
                self._parse_copywriting_rules(rules, data["copywriting"])

            # 解析布局（兼容旧格式）
            if data.get("layout"):
                self._parse_layout_rules(rules, data["layout"])

            # 解析规则列表：新格式用 "rules"，旧格式用 "secondary_rules"
            raw_rules = data.get("rules") or data.get("secondary_rules") or []
            if raw_rules:
                self._parse_secondary_rules(rules, raw_rules)

            # 解析前置条件
            if data.get("preconditions"):
                rules.preconditions = data["preconditions"]

            logger.info(f"LLM规则提取完成: brand_name={rules.brand_name}, "
                       f"secondary_rules={len(rules.secondary_rules)}项")

            return rules

        except Exception as e:
            logger.error(f"LLM规则提取失败: {type(e).__name__}: {e}")
            import traceback
            logger.error(traceback.format_exc())
            return BrandRules(
                brand_id="",
                brand_name="",
                version="1.0",
                source=filename,
            )

    def _extract_rules_with_llm_stream(self, text: str, filename: str, stream_callback=None):
        """
        使用LLM流式提取品牌规范

        Args:
            text: 文档文本内容
            filename: 文件名
            stream_callback: 流式回调函数，接收每个文本块

        Yields:
            每个文本块
        """
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI

        logger.info(f"使用LLM流式提取品牌规范: {filename}")

        # 检查API配置
        if not settings.llm_api_key:
            logger.warning("未配置文本模型 API Key，无法进行LLM提取")
            if stream_callback:
                stream_callback("[错误] 未配置文本模型 API Key")
            return BrandRules(
                brand_id="",
                brand_name="",
                version="1.0",
                source=filename,
            )

        system_prompt = PARSE_SYSTEM_PROMPT

        # 截取文本
        max_chars = 60000
        if len(text) > max_chars:
            text_for_llm = text[:max_chars]
            logger.info(f"文本过长({len(text)}字符)，截取前{max_chars}字符")
        else:
            text_for_llm = text

        user_prompt = f"""请从以下品牌规范文档中提取结构化规则。

文档内容：
{text_for_llm}

请输出JSON格式的规则："""

        full_content = ""

        try:
            logger.info(f"文本模型配置: base={settings.llm_api_base}, model={settings.llm_model}")

            llm = ChatOpenAI(
                model=settings.llm_model,
                base_url=settings.llm_api_base,
                api_key=settings.llm_api_key,
                temperature=0,
                timeout=300,
                max_tokens=16384,  # 规则文档解析输出量大，必须显式设置
            )

            messages = [
                SystemMessage(content=system_prompt),
                HumanMessage(content=user_prompt),
            ]

            logger.info("正在调用LLM流式API...")

            for chunk in llm.stream(messages):
                if chunk.content:
                    text_chunk = chunk.content
                    full_content += text_chunk

                    if stream_callback:
                        stream_callback(text_chunk)

                    yield text_chunk

        except Exception as e:
            logger.error(f"LLM流式规则提取失败: {type(e).__name__}: {e}")
            import traceback
            logger.error(traceback.format_exc())
            if stream_callback:
                stream_callback(f"\n\n[错误] 提取失败: {str(e)}")

    def parse_stream_result(self, full_content: str, filename: str) -> BrandRules:
        """
        解析流式输出的完整结果

        Args:
            full_content: 流式输出的完整文本
            filename: 文件名

        Returns:
            BrandRules对象
        """
        data = self._parse_json_response(full_content)

        if data is None:
            logger.error("无法从LLM响应中解析JSON")
            try:
                import time
                dump_path = Path(__file__).parent.parent.parent / "docs" / f"llm_response_dump_{int(time.time())}.txt"
                dump_path.parent.mkdir(parents=True, exist_ok=True)
                dump_path.write_text(full_content, encoding="utf-8")
                logger.error(f"完整LLM响应已写入: {dump_path}  (共{len(full_content)}字符)")
            except Exception:
                pass
            return BrandRules(
                brand_id="",
                brand_name="",
                version="1.0",
                source=filename,
            )

        # 构建BrandRules
        rules = BrandRules(
            brand_id="",
            brand_name=data.get("brand_name") or "",
            version="1.0",
            source=filename,
        )

        # 解析各项规则（兼容旧格式）
        if data.get("color"):
            self._parse_color_rules(rules, data["color"])
        if data.get("logo"):
            self._parse_logo_rules(rules, data["logo"])
        if data.get("font"):
            self._parse_font_rules(rules, data["font"])
        if data.get("copywriting"):
            self._parse_copywriting_rules(rules, data["copywriting"])
        if data.get("layout"):
            self._parse_layout_rules(rules, data["layout"])
        # 新格式用 "rules"，旧格式用 "secondary_rules"
        raw_rules = data.get("rules") or data.get("secondary_rules") or []
        if raw_rules:
            self._parse_secondary_rules(rules, raw_rules)
        if data.get("preconditions"):
            rules.preconditions = data["preconditions"]

        return rules

    def _parse_json_response(self, content: str) -> dict | None:
        """从LLM响应中解析JSON（委托给公共方法）"""
        result = parse_json_response(content)
        return result if isinstance(result, dict) else None

    def _parse_color_rules(self, rules: BrandRules, color_data: dict):
        """解析色彩规则"""
        rules.color = ColorRules()

        if color_data.get("primary"):
            primary = color_data["primary"]
            value = primary.get("value") or primary.get("hex") or ""
            if value:
                rules.color.primary = ColorRule(
                    name=primary.get("name", "主色"),
                    value=value,
                )

        if color_data.get("secondary"):
            secondary_list = [
                ColorRule(name=c.get("name", ""), value=c.get("value") or c.get("hex") or "")
                for c in color_data["secondary"]
                if c.get("value") or c.get("hex")
            ]
            if secondary_list:
                rules.color.secondary = secondary_list

        if color_data.get("forbidden"):
            forbidden_list = [
                ColorRule(
                    name=c.get("name", ""),
                    value=c.get("value") or c.get("hex") or "",
                    reason=c.get("reason", ""),
                )
                for c in color_data["forbidden"]
                if c.get("value") or c.get("hex")
            ]
            if forbidden_list:
                rules.color.forbidden = forbidden_list

        # 解析额外的色彩规则
        if color_data.get("additional_rules"):
            rules.color.additional_rules = [r for r in color_data["additional_rules"] if r]

        if color_data.get("description"):
            rules.color.description = color_data["description"]

    def _parse_logo_rules(self, rules: BrandRules, logo_data: dict):
        """解析Logo规则"""
        size_range = logo_data.get("size_range")

        # 只在 LLM 明确返回了 size_range 时才解析，避免注入虚假默认值
        if size_range is None:
            parsed_size_range = None
        elif isinstance(size_range, str):
            numbers = re.findall(r'(\d+\.?\d*)', size_range)
            parsed_size_range = {
                "min": int(float(numbers[0])) if numbers else None,
                "max": int(float(numbers[1])) if len(numbers) > 1 else None,
            }
        elif isinstance(size_range, dict):
            min_val = size_range.get("min")
            max_val = size_range.get("max")
            parsed_size_range = {
                "min": int(float(min_val)) if min_val is not None else None,
                "max": int(float(max_val)) if max_val is not None else None,
            }
        else:
            parsed_size_range = None

        # safe_margin_px：只在 LLM 明确返回时才设置
        raw_margin = logo_data.get("safe_margin_px")
        safe_margin = int(float(raw_margin)) if raw_margin is not None else None

        rules.logo = LogoRules(
            position=logo_data.get("position", "") or "",
            position_description=logo_data.get("position_description", "") or "",
            size_range=parsed_size_range,
            safe_margin_px=safe_margin,
            additional_rules=[r for r in (logo_data.get("additional_rules") or []) if r],
            min_display_ratio=logo_data.get("min_display_ratio") or logo_data.get("min_size_ratio"),
            color_requirements=[r for r in (logo_data.get("color_requirements") or []) if r],
            background_requirements=[r for r in (logo_data.get("background_requirements") or []) if r],
        )

    def _parse_font_rules(self, rules: BrandRules, font_data: dict):
        """解析字体规则"""
        allowed = font_data.get("allowed") or []
        forbidden = font_data.get("forbidden") or []
        size_rules = font_data.get("size_rules") or {}

        # 过滤空值
        allowed = [f for f in allowed if f]
        forbidden = [f for f in forbidden if f]

        if allowed or forbidden or size_rules:
            rules.font = FontRules(
                allowed=allowed,
                forbidden=forbidden,
                size_rules=size_rules,
                additional_rules=[r for r in (font_data.get("additional_rules") or []) if r],
                note=font_data.get("note") or None,
            )

    def _parse_copywriting_rules(self, rules: BrandRules, cw_data: dict):
        """解析文案规则"""
        forbidden_words = cw_data.get("forbidden_words") or []

        if forbidden_words:
            rules.copywriting = CopywritingRules()
            rules.copywriting.forbidden_words = [
                ForbiddenWord(word=w.get("word", ""), category=w.get("category") or "禁用词")
                for w in forbidden_words
                if w.get("word")
            ]

        required_content = cw_data.get("required_content") or []
        if required_content and rules.copywriting:
            rules.copywriting.required_content = required_content

    def _parse_layout_rules(self, rules: BrandRules, layout_data: dict):
        """解析布局规则"""
        margin_min = layout_data.get("margin_min")
        description = layout_data.get("description") or ""

        if margin_min is not None or description:
            rules.layout = LayoutRules(
                margin_min=int(margin_min) if margin_min else 20,
                description=description,
            )

    def _parse_secondary_rules(self, rules: BrandRules, secondary_data: list):
        """解析次要规范"""
        for item in secondary_data:
            if isinstance(item, dict) and item.get("name") and item.get("content"):
                rule = SecondaryRule(
                    category=item.get("category", "其他"),
                    name=item.get("name", ""),
                    content=item.get("content", ""),
                    priority=item.get("priority", 1),
                    rule_source_id=item.get("rule_source_id") or None,
                    # 新三段式字段
                    fail_condition=item.get("fail_condition") or None,
                    review_condition=item.get("review_condition") or None,
                    pass_condition=item.get("pass_condition") or None,
                    # 旧字段兼容
                    output_level=item.get("output_level") or None,
                    threshold=item.get("threshold") or None,
                    feedback_text=item.get("feedback_text") or None,
                )
                rules.secondary_rules.append(rule)

    def extract_text_only(self, file_data: bytes, filename: str) -> str:
        """仅提取文档文本，不调用LLM解析（用于多文件合并场景）"""
        ext = Path(filename).suffix.lower()

        if ext not in SUPPORTED_FORMATS:
            raise ValueError(f"不支持的文档格式: {ext}")

        method_name, _ = SUPPORTED_FORMATS[ext]
        parse_method = getattr(self, method_name)
        file_path = Path(filename)

        # 根据文件类型提取文本
        if ext == ".pdf":
            doc = fitz.open(stream=file_data, filetype="pdf")
            text_content = []
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    text_content.append(f"=== 第{page_num + 1}页 ===\n{text}")
            doc.close()
            return "\n\n".join(text_content)

        elif ext in (".ppt", ".pptx"):
            prs = Presentation(io.BytesIO(file_data))
            text_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text_content.append(f"=== 第{slide_num}页 ===\n" + "\n".join(slide_texts))
            return "\n\n".join(text_content)

        elif ext in (".doc", ".docx"):
            text_content = []
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_data))
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_content.append(para.text.strip())
                for table_idx, table in enumerate(doc.tables):
                    table_texts = []
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            table_texts.append(row_text)
                    if table_texts:
                        text_content.append(f"\n=== 表格{table_idx + 1} ===\n" + "\n".join(table_texts))
            except Exception as e:
                logger.warning(f"解析Word失败: {e}")
                text_content.append(file_data.decode('utf-8', errors='ignore'))
            return "\n\n".join(text_content)

        elif ext in (".xls", ".xlsx"):
            text_content = []
            try:
                if ext == ".xlsx":
                    from openpyxl import load_workbook
                    wb = load_workbook(io.BytesIO(file_data), data_only=True)
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        sheet_texts = [f"=== 工作表: {sheet_name} ==="]
                        for row in sheet.iter_rows(values_only=True):
                            row_values = [str(cell) if cell is not None else "" for cell in row]
                            row_text = " | ".join(row_values)
                            if row_text.strip(" |"):
                                sheet_texts.append(row_text)
                        if len(sheet_texts) > 1:
                            text_content.append("\n".join(sheet_texts))
                else:
                    import xlrd
                    workbook = xlrd.open_workbook(file_contents=file_data)
                    for sheet in workbook.sheets():
                        sheet_texts = [f"=== 工作表: {sheet.name} ==="]
                        for row_idx in range(sheet.nrows):
                            row_values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                            row_text = " | ".join(row_values)
                            if row_text.strip(" |"):
                                sheet_texts.append(row_text)
                        if len(sheet_texts) > 1:
                            text_content.append("\n".join(sheet_texts))
            except Exception as e:
                logger.warning(f"解析Excel失败: {e}")
            return "\n\n".join(text_content)

        elif ext in (".md", ".txt"):
            encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
            for encoding in encodings:
                try:
                    return file_data.decode(encoding)
                except UnicodeDecodeError:
                    continue
            return file_data.decode('utf-8', errors='ignore')

        return ""


    # ── 异步包装方法 ──────────────────────────────────────────────────────────────

    async def async_parse(self, file_data: bytes, filename: str) -> "BrandRules":
        """parse 的异步版本，避免阻塞事件循环（LLM HTTP 调用是同步阻塞的）"""
        import asyncio
        return await asyncio.to_thread(self.parse, file_data, filename)

    async def async_extract_rules_with_llm(self, text: str, filename: str) -> "BrandRules":
        """_extract_rules_with_llm 的异步版本"""
        import asyncio
        return await asyncio.to_thread(self._extract_rules_with_llm, text, filename)

    async def async_extract_text_only(self, file_data: bytes, filename: str) -> str:
        """extract_text_only 的异步版本（PDF/PPT 解析也可能耗时）"""
        import asyncio
        return await asyncio.to_thread(self.extract_text_only, file_data, filename)


# 全局文档解析器实例
document_parser = DocumentParser()
